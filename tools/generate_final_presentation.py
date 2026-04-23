from __future__ import annotations

from pathlib import Path
from textwrap import wrap

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSET_DIR = ROOT / "presentation_assets"
OUTPUT = ROOT / "Trend-Fx_Final_Presentation.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

NAVY = RGBColor(12, 26, 58)
DEEP = RGBColor(20, 42, 82)
TEAL = RGBColor(0, 153, 153)
SKY = RGBColor(96, 160, 255)
MINT = RGBColor(81, 193, 149)
GOLD = RGBColor(240, 181, 62)
RED = RGBColor(214, 85, 85)
WHITE = RGBColor(255, 255, 255)
INK = RGBColor(41, 51, 77)
MUTED = RGBColor(104, 117, 142)
PANEL = RGBColor(245, 247, 251)
LINE = RGBColor(212, 220, 232)


def ensure_assets():
    ASSET_DIR.mkdir(exist_ok=True)


def font(size: int, bold: bool = False):
    candidates = [
        "C:/Windows/Fonts/arialbd.ttf" if bold else "C:/Windows/Fonts/arial.ttf",
        "C:/Windows/Fonts/segoeuib.ttf" if bold else "C:/Windows/Fonts/segoeui.ttf",
    ]
    for candidate in candidates:
        path = Path(candidate)
        if path.exists():
            return ImageFont.truetype(str(path), size=size)
    return ImageFont.load_default()


def draw_wrapped(draw, text, xy, max_width, fill, body_font, line_gap=8):
    x, y = xy
    words = text.split()
    lines = []
    current = ""
    for word in words:
        test = f"{current} {word}".strip()
        if draw.textlength(test, font=body_font) <= max_width:
            current = test
        else:
            if current:
                lines.append(current)
            current = word
    if current:
        lines.append(current)

    for line in lines:
        draw.text((x, y), line, font=body_font, fill=fill)
        y += body_font.size + line_gap
    return y


def make_hero_image(path: Path):
    img = Image.new("RGB", (1600, 900), (10, 24, 54))
    draw = ImageDraw.Draw(img)

    for i in range(12):
        draw.rounded_rectangle((60 + i * 110, 560 - i * 20, 130 + i * 110, 780 - i * 12), radius=14, fill=(24, 56, 102))

    for x in range(0, 1600, 80):
        draw.line((x, 0, x, 900), fill=(25, 54, 92), width=1)
    for y in range(0, 900, 80):
        draw.line((0, y, 1600, y), fill=(25, 54, 92), width=1)

    title_font = font(56, bold=True)
    sub_font = font(28)
    small_font = font(22)
    draw.text((90, 92), "Trend-Fx", font=title_font, fill=(255, 255, 255))
    draw.text((92, 168), "Forex intelligence dashboard for faster retail decision-making", font=sub_font, fill=(171, 199, 244))
    draw.text((92, 230), "Professional presentation visual", font=small_font, fill=(108, 224, 202))

    panel = (840, 120, 1490, 760)
    draw.rounded_rectangle(panel, radius=28, fill=(247, 249, 252))
    draw.text((890, 170), "Live Market Snapshot", font=font(34, bold=True), fill=(34, 48, 78))

    chart_box = (890, 240, 1440, 540)
    draw.rounded_rectangle(chart_box, radius=20, outline=(214, 222, 235), width=3, fill=(255, 255, 255))
    for y in range(270, 520, 52):
        draw.line((920, y, 1410, y), fill=(228, 233, 241), width=2)
    for x in range(930, 1410, 64):
        draw.line((x, 270, x, 510), fill=(235, 239, 246), width=1)

    points = [(930, 470), (1000, 430), (1060, 445), (1130, 390), (1200, 350), (1270, 315), (1340, 280), (1400, 305)]
    draw.line(points, fill=(26, 120, 255), width=7)
    for x, y in points:
        draw.ellipse((x - 8, y - 8, x + 8, y + 8), fill=(26, 120, 255))

    cards = [
        ("Pair", "EUR/USD", (890, 575, 1060, 710), (236, 246, 255)),
        ("Signal", "BUY", (1080, 575, 1250, 710), (232, 249, 241)),
        ("Confidence", "78%", (1270, 575, 1440, 710), (255, 247, 230)),
    ]
    for label, value, box, fill in cards:
        draw.rounded_rectangle(box, radius=18, fill=fill, outline=(215, 224, 237), width=2)
        draw.text((box[0] + 22, box[1] + 20), label, font=font(22), fill=(102, 115, 138))
        draw.text((box[0] + 22, box[1] + 60), value, font=font(34, bold=True), fill=(28, 43, 71))

    img.save(path)


def make_architecture_diagram(path: Path):
    img = Image.new("RGB", (1600, 900), (251, 252, 254))
    draw = ImageDraw.Draw(img)
    title_font = font(42, bold=True)
    sub_font = font(24)
    box_title = font(28, bold=True)
    box_body = font(20)
    draw.text((70, 40), "Trend-Fx System Architecture", font=title_font, fill=(18, 35, 68))
    draw.text((70, 95), "Client-driven dashboard backed by FastAPI services and SQLite persistence", font=sub_font, fill=(98, 112, 136))

    boxes = {
        "users": ((80, 220, 330, 470), "End Users", "Students and retail forex learners\nuse dashboard, converter,\nauth, and market insights", (230, 240, 255)),
        "frontend": ((420, 180, 760, 520), "React + Vite Frontend", "Routes: Home, About, Converter,\nDashboard, Login/Register,\nForgot/Reset Password", (237, 246, 255)),
        "backend": ((850, 180, 1190, 520), "FastAPI Backend", "Auth router, users router,\nmarket/news services,\nJWT validation and password reset", (232, 247, 241)),
        "db": ((1280, 180, 1510, 520), "SQLite Database", "users\npassword_reset_tokens", (255, 246, 229)),
        "apis": ((850, 610, 1190, 810), "External Data Services", "Twelve Data API for live forex prices\nSMTP server for reset-code emails", (248, 239, 255)),
        "devops": ((420, 610, 760, 810), "GitHub Workflow", "Feature branches, pull requests,\ncode reviews, sprint planning,\nissue tracking", (245, 247, 250)),
    }

    for box in boxes.values():
        x1, y1, x2, y2 = box[0]
        draw.rounded_rectangle((x1, y1, x2, y2), radius=28, fill=box[3], outline=(214, 222, 235), width=3)
        draw.text((x1 + 24, y1 + 22), box[1], font=box_title, fill=(28, 42, 70))
        draw_wrapped(draw, box[2], (x1 + 24, y1 + 80), x2 - x1 - 48, (92, 106, 130), box_body, line_gap=6)

    arrows = [
        ((330, 345), (420, 345)),
        ((760, 345), (850, 345)),
        ((1190, 345), (1280, 345)),
        ((1000, 520), (1000, 610)),
        ((760, 710), (850, 710)),
        ((590, 520), (590, 610)),
    ]
    for start, end in arrows:
        draw.line((start, end), fill=(42, 78, 134), width=7)
        ex, ey = end
        draw.polygon([(ex, ey), (ex - 18, ey - 12), (ex - 18, ey + 12)], fill=(42, 78, 134))

    img.save(path)


def make_erd_diagram(path: Path):
    img = Image.new("RGB", (1600, 900), (250, 252, 255))
    draw = ImageDraw.Draw(img)
    draw.text((70, 50), "Database ERD", font=font(42, True), fill=(19, 35, 67))
    draw.text((70, 105), "Authentication data is persisted in a compact relational schema", font=font(24), fill=(99, 113, 136))

    def table(box, title, fields, fill):
        x1, y1, x2, y2 = box
        draw.rounded_rectangle(box, radius=24, fill=fill, outline=(207, 217, 230), width=3)
        draw.rounded_rectangle((x1, y1, x2, y1 + 72), radius=24, fill=(24, 54, 100))
        draw.rectangle((x1, y1 + 48, x2, y1 + 72), fill=(24, 54, 100))
        draw.text((x1 + 24, y1 + 18), title, font=font(30, True), fill=(255, 255, 255))
        current_y = y1 + 104
        for field in fields:
            draw.text((x1 + 28, current_y), field, font=font(23), fill=(50, 64, 92))
            current_y += 52

    table((190, 220, 680, 670), "users", [
        "PK  id : integer",
        "email : varchar(255) unique",
        "full_name : varchar(255)",
        "password_hash : varchar(255)",
        "created_at : datetime",
    ], (236, 244, 255))

    table((920, 220, 1410, 670), "password_reset_tokens", [
        "PK  id : integer",
        "FK  user_id -> users.id",
        "token : varchar(255) unique",
        "is_used : boolean",
        "expires_at : datetime",
        "created_at : datetime",
    ], (236, 248, 241))

    draw.line((680, 450, 920, 450), fill=(47, 88, 146), width=8)
    draw.polygon([(920, 450), (895, 435), (895, 465)], fill=(47, 88, 146))
    draw.text((740, 395), "1 user", font=font(22, True), fill=(47, 88, 146))
    draw.text((742, 455), "many reset codes", font=font(22), fill=(85, 102, 129))

    note_box = (330, 730, 1270, 835)
    draw.rounded_rectangle(note_box, radius=22, fill=(255, 247, 230), outline=(237, 211, 156), width=2)
    draw.text((365, 758), "Operational market/news data is fetched live from APIs and not persisted in relational tables.", font=font(24), fill=(115, 84, 27))

    img.save(path)


def make_frontend_structure(path: Path):
    img = Image.new("RGB", (1600, 900), (251, 252, 254))
    draw = ImageDraw.Draw(img)
    draw.text((70, 50), "Frontend Component Structure", font=font(42, True), fill=(18, 35, 68))
    draw.text((70, 105), "React pages, shared layout components, and service modules organized for clear separation of concerns", font=font(24), fill=(98, 112, 136))

    tree = [
        ("App.jsx", (660, 150, 940, 230), (235, 244, 255)),
        ("Pages", (180, 320, 520, 760), (242, 248, 255)),
        ("Shared Components", (630, 320, 970, 760), (236, 248, 241)),
        ("Service Layer", (1080, 320, 1420, 760), (255, 247, 231)),
    ]
    for label, box, fill in tree:
        draw.rounded_rectangle(box, radius=24, fill=fill, outline=(211, 220, 233), width=3)
        draw.text((box[0] + 24, box[1] + 20), label, font=font(30, True), fill=(29, 44, 72))

    pages = ["Home", "About", "Converter", "Dashboard", "Login", "Register", "ForgotPassword", "ResetPassword", "Logout"]
    comps = ["Navbar", "HeroLayout", "MarketChart", "Clock", "RightSidebar", "TopTabs", "ProfileCard", "ProtectedRoute", "PublicNavbar"]
    services = ["api.js", "auth.js", "authApi.js", "market.js", "news.js"]

    def add_list(items, start_x, start_y, width, fill):
        y = start_y
        for item in items:
            draw.rounded_rectangle((start_x, y, start_x + width, y + 42), radius=12, fill=fill, outline=(214, 222, 235), width=2)
            draw.text((start_x + 14, y + 8), item, font=font(20), fill=(56, 70, 96))
            y += 52

    add_list(pages, 210, 380, 280, (255, 255, 255))
    add_list(comps, 660, 380, 280, (255, 255, 255))
    add_list(services, 1110, 380, 280, (255, 255, 255))

    for x in (350, 800, 1250):
        draw.line((800, 230, x, 320), fill=(47, 88, 146), width=7)
        draw.polygon([(x, 320), (x - 18, 306), (x + 18, 306)], fill=(47, 88, 146))

    img.save(path)


def make_project_board(path: Path):
    img = Image.new("RGB", (1600, 900), (248, 250, 253))
    draw = ImageDraw.Draw(img)
    draw.text((70, 50), "Final Project Board Snapshot", font=font(42, True), fill=(18, 35, 68))
    draw.text((70, 105), "Kanban-style summary reconstructed from branches, commits, and merged pull requests", font=font(24), fill=(98, 112, 136))

    columns = [
        ("Completed", (80, 180, 500, 820), (231, 247, 239)),
        ("Validated / Merged", (590, 180, 1010, 820), (235, 244, 255)),
        ("Future Backlog", (1100, 180, 1520, 820), (255, 247, 231)),
    ]
    cards = {
        "Completed": [
            ("Frontend design + Vite setup", "Preet", "Sprint 1"),
            ("FastAPI backend + JWT auth", "Vignesh", "Sprint 2"),
            ("Live price tracking API", "Vignesh", "Sprint 2"),
            ("RSI + MACD dashboard charts", "Preet", "Sprint 3"),
            ("Auth flow, reset, logout routing", "Vignesh", "Sprint 3"),
            ("About page and bug fixes", "Preet", "Sprint 3"),
        ],
        "Validated / Merged": [
            ("Branch: features-preet", "Merged PR #2 / #7", "Done"),
            ("Branch: features-vignesh", "Merged PR #3 / #4 / #5", "Done"),
            ("Error handling hardening", "Vignesh", "Done"),
            ("Navigation + API integration", "Preet", "Done"),
            ("Testing and cleanup", "Team", "Done"),
        ],
        "Future Backlog": [
            ("Production email SMTP secrets", "Team", "Pending"),
            ("Backtesting module", "Team", "Future"),
            ("Risk management insights", "Team", "Future"),
            ("Machine learning experimentation", "Team", "Future"),
        ],
    }

    for title, box, fill in columns:
        draw.rounded_rectangle(box, radius=28, fill=fill, outline=(210, 220, 232), width=3)
        draw.text((box[0] + 24, box[1] + 22), title, font=font(30, True), fill=(29, 44, 72))
        y = box[1] + 80
        for card_title, owner, status in cards[title]:
            card = (box[0] + 20, y, box[2] - 20, y + 98)
            draw.rounded_rectangle(card, radius=18, fill=(255, 255, 255), outline=(216, 223, 235), width=2)
            draw.text((card[0] + 18, card[1] + 16), card_title, font=font(22, True), fill=(48, 63, 90))
            draw.text((card[0] + 18, card[1] + 54), owner, font=font(18), fill=(84, 99, 124))
            badge_w = 110
            draw.rounded_rectangle((card[2] - badge_w - 18, card[1] + 18, card[2] - 18, card[1] + 52), radius=16, fill=(230, 239, 251))
            draw.text((card[2] - badge_w, card[1] + 24), status, font=font(17, True), fill=(42, 78, 134))
            y += 114

    img.save(path)


def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, size=20, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    box.text_frame.word_wrap = True
    return box


def add_bullets(slide, left, top, width, height, items, color=INK, size=20, bullet_color=TEAL):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(6)
        run = p.runs[0]
        run.font.size = Pt(size)
        run.font.color.rgb = color
        p.bullet = True
    return box


def add_title(slide, title, subtitle=None, dark=True):
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(6.8), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    run = p.runs[0]
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = WHITE if dark else NAVY
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.72), Inches(1.08), Inches(9.5), Inches(0.55))
        tf2 = sub.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        r2 = p2.runs[0]
        r2.font.size = Pt(13)
        r2.font.color.rgb = RGBColor(214, 225, 246) if dark else MUTED


def add_footer(slide, page_no):
    add_textbox(slide, Inches(11.5), Inches(7.0), Inches(1.0), Inches(0.25), str(page_no), size=11, color=MUTED, align=PP_ALIGN.RIGHT)


def build_presentation():
    ensure_assets()
    hero = ASSET_DIR / "hero_forex.png"
    architecture = ASSET_DIR / "architecture_diagram.png"
    erd = ASSET_DIR / "erd_diagram.png"
    frontend = ASSET_DIR / "frontend_structure.png"
    board = ASSET_DIR / "project_board.png"

    make_hero_image(hero)
    make_architecture_diagram(architecture)
    make_erd_diagram(erd)
    make_frontend_structure(frontend)
    make_project_board(board)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    slide.shapes.add_picture(str(hero), Inches(6.7), Inches(0.5), width=Inches(6.0))
    add_textbox(slide, Inches(0.75), Inches(0.9), Inches(5.5), Inches(1.2), "Trend-Fx", size=30, color=WHITE, bold=True)
    add_textbox(slide, Inches(0.75), Inches(1.75), Inches(5.4), Inches(1.3), "Final Project Presentation", size=26, color=RGBColor(172, 198, 246), bold=True)
    add_textbox(slide, Inches(0.75), Inches(3.0), Inches(5.4), Inches(1.8), "Forex market prediction dashboard with live data, authentication, analytics, and project-managed delivery.", size=18, color=RGBColor(224, 233, 248))
    add_textbox(slide, Inches(0.75), Inches(5.9), Inches(5.6), Inches(0.8), "Team: Vignesh + Preet", size=18, color=RGBColor(110, 227, 206), bold=True)
    add_textbox(slide, Inches(0.75), Inches(6.35), Inches(5.6), Inches(0.6), "Final review deck", size=14, color=RGBColor(190, 202, 227))
    add_footer(slide, 1)

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_title(slide, "1. Project Introduction", "Problem statement, target users, goals, and unique value proposition", dark=False)
    add_textbox(slide, Inches(0.7), Inches(1.7), Inches(2.8), Inches(0.4), "Problem Statement", size=20, color=NAVY, bold=True)
    add_bullets(slide, Inches(0.7), Inches(2.1), Inches(5.6), Inches(1.4), [
        "Retail forex learners struggle to interpret fast-moving market data and technical indicators in one place.",
        "Most beginner workflows depend on multiple fragmented tools for price checks, signals, and account access.",
    ], size=17)
    add_textbox(slide, Inches(0.7), Inches(3.65), Inches(2.8), Inches(0.4), "Target Users", size=20, color=NAVY, bold=True)
    add_bullets(slide, Inches(0.7), Inches(4.0), Inches(5.6), Inches(1.1), [
        "Students learning forex concepts",
        "Retail traders who need a lightweight dashboard",
        "Project evaluators looking for a complete full-stack workflow",
    ], size=17)
    add_textbox(slide, Inches(6.7), Inches(1.7), Inches(2.8), Inches(0.4), "Project Goals", size=20, color=NAVY, bold=True)
    add_bullets(slide, Inches(6.7), Inches(2.1), Inches(5.8), Inches(1.6), [
        "Deliver live forex visibility with RSI and MACD-based insights",
        "Support secure login, registration, reset, and logout flows",
        "Provide an intuitive React dashboard backed by FastAPI services",
    ], size=17)
    add_textbox(slide, Inches(6.7), Inches(4.1), Inches(3.5), Inches(0.4), "Unique Value Proposition", size=20, color=NAVY, bold=True)
    quote = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(4.5), Inches(5.6), Inches(1.5))
    quote.fill.solid()
    quote.fill.fore_color.rgb = PANEL
    quote.line.color.rgb = LINE
    add_textbox(slide, Inches(6.95), Inches(4.78), Inches(5.1), Inches(0.95), "Trend-Fx combines authentication, market data, indicators, and a learning-oriented dashboard in one focused academic product.", size=18, color=INK)
    add_footer(slide, 2)

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(246, 248, 252))
    add_title(slide, "2. Technical Overview", "End-to-end stack, integrations, and engineering choices", dark=False)
    sections = [
        ("Frontend", ["React 19", "Vite 7", "React Router DOM", "Recharts", "Custom CSS"]),
        ("Backend", ["FastAPI", "Uvicorn", "SQLAlchemy", "JWT auth", "Passlib password hashing"]),
        ("Database", ["SQLite test.db", "users table", "password_reset_tokens table"]),
        ("Tools", ["Git + GitHub", "Feature branches", "Pull requests", "VS Code", "PowerShell"]),
        ("Integrations", ["Twelve Data API", "SMTP email flow", "dotenv environment config"]),
    ]
    x_positions = [0.7, 3.15, 5.6, 8.05, 10.5]
    widths = [2.1, 2.1, 2.1, 2.1, 2.1]
    fills = [RGBColor(236, 244, 255), RGBColor(236, 248, 241), RGBColor(255, 247, 231), RGBColor(246, 241, 255), RGBColor(240, 247, 249)]
    for (title, items), x, w, fill in zip(sections, x_positions, widths, fills):
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.95), Inches(w), Inches(4.5))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.color.rgb = LINE
        add_textbox(slide, Inches(x + 0.18), Inches(2.15), Inches(w - 0.36), Inches(0.35), title, size=18, color=NAVY, bold=True)
        add_bullets(slide, Inches(x + 0.16), Inches(2.55), Inches(w - 0.32), Inches(3.6), items, size=14)
    add_footer(slide, 3)

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_title(slide, "3. System Architecture", "How data, authentication, and UI layers interact", dark=False)
    slide.shapes.add_picture(str(architecture), Inches(0.75), Inches(1.55), width=Inches(11.85))
    add_footer(slide, 4)

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(246, 248, 252))
    add_title(slide, "4. Backend API Overview", "Core routes powering auth, user state, market data, and news", dark=False)
    api_cards = [
        ("Auth", ["/api/auth/register", "/api/auth/login", "/api/auth/forgot-password", "/api/auth/reset-password"], RGBColor(236, 244, 255)),
        ("User", ["/api/users/me"], RGBColor(236, 248, 241)),
        ("Market", ["/market-data", "/price", "/api/market/live", "/api/market/candles"], RGBColor(255, 247, 231)),
        ("News", ["/api/news/latest"], RGBColor(246, 241, 255)),
    ]
    positions = [(0.75, 1.85), (3.85, 1.85), (6.95, 1.85), (10.05, 1.85)]
    for (title, items, fill), (x, y) in zip(api_cards, positions):
        shp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(2.45), Inches(3.45))
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        shp.line.color.rgb = LINE
        add_textbox(slide, Inches(x + 0.18), Inches(y + 0.2), Inches(2.1), Inches(0.3), title, size=18, color=NAVY, bold=True)
        add_bullets(slide, Inches(x + 0.15), Inches(y + 0.55), Inches(2.15), Inches(2.6), items, size=14)
    quote = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(5.6), Inches(11.8), Inches(1.0))
    quote.fill.solid()
    quote.fill.fore_color.rgb = WHITE
    quote.line.color.rgb = LINE
    add_textbox(slide, Inches(1.0), Inches(5.9), Inches(11.2), Inches(0.4), "API design separates concerns cleanly: authentication and user state are secured, while market/news routes support analytical and UI needs.", size=17, color=INK, align=PP_ALIGN.CENTER)
    add_footer(slide, 5)

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_title(slide, "5. Database Schema", "Entity relationship diagram for authentication and password recovery", dark=False)
    slide.shapes.add_picture(str(erd), Inches(0.7), Inches(1.55), width=Inches(11.9))
    add_footer(slide, 6)

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(246, 248, 252))
    add_title(slide, "6. Frontend Structure", "Pages, shared components, and service modules", dark=False)
    slide.shapes.add_picture(str(frontend), Inches(0.72), Inches(1.55), width=Inches(11.85))
    add_footer(slide, 7)

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_title(slide, "7. Team Workflow & Project Management", "Branching, merged pull requests, and a final project board view", dark=False)
    slide.shapes.add_picture(str(board), Inches(0.6), Inches(1.5), width=Inches(12.1))
    add_footer(slide, 8)

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(246, 248, 252))
    add_title(slide, "8. Coordination Reflection", "How planning and version control improved team delivery", dark=False)
    left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(5.8), Inches(4.7))
    left.fill.solid()
    left.fill.fore_color.rgb = WHITE
    left.line.color.rgb = LINE
    add_textbox(slide, Inches(1.05), Inches(2.05), Inches(5.2), Inches(0.4), "Evidence from Repository Workflow", size=20, color=NAVY, bold=True)
    add_bullets(slide, Inches(1.02), Inches(2.45), Inches(5.15), Inches(3.6), [
        "Feature branches: features-preet and features-vignesh",
        "Merged pull requests observed in history: #2, #3, #4, #5, #7",
        "Commit split shows Preet focused on UI, charts, and navigation while Vignesh focused on backend, APIs, auth, and stabilization",
        "Issue-style iteration visible through repeated bug-fix and enhancement commits",
    ], size=17)
    right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.7))
    right.fill.solid()
    right.fill.fore_color.rgb = WHITE
    right.line.color.rgb = LINE
    add_textbox(slide, Inches(7.15), Inches(2.05), Inches(5.0), Inches(0.4), "Reflection on Board Usage", size=20, color=NAVY, bold=True)
    add_bullets(slide, Inches(7.12), Inches(2.45), Inches(4.95), Inches(3.6), [
        "Clear task ownership reduced duplicate work across frontend and backend streams",
        "Sprint grouping helped the team move from foundation -> APIs -> analytics -> polish",
        "Branches and PRs made integration safer and easier to review before merging",
        "A visible board improved communication on status, blockers, and completion readiness",
    ], size=17)
    add_footer(slide, 9)

    # Slide 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    add_title(slide, "9. Closing Summary", "What the final project demonstrates", dark=True)
    add_bullets(slide, Inches(0.95), Inches(1.8), Inches(5.6), Inches(3.7), [
        "A complete full-stack academic product from UI to API to database",
        "Live forex data visualization with technical indicators",
        "Secure user lifecycle: register, login, logout, reset-password",
        "Team delivery supported by branches, PRs, and sprint-based coordination",
    ], color=WHITE, size=20)
    summary = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.8), Inches(5.0), Inches(3.3))
    summary.fill.solid()
    summary.fill.fore_color.rgb = RGBColor(24, 44, 82)
    summary.line.color.rgb = RGBColor(65, 98, 151)
    add_textbox(slide, Inches(7.3), Inches(2.15), Inches(4.4), Inches(1.8), "Trend-Fx turns fragmented forex monitoring into a guided dashboard experience that is easier to learn, safer to access, and faster to evaluate.", size=22, color=WHITE)
    add_textbox(slide, Inches(7.3), Inches(4.3), Inches(4.4), Inches(0.4), "Thank you", size=24, color=RGBColor(110, 227, 206), bold=True)
    add_footer(slide, 10)

    prs.save(OUTPUT)


if __name__ == "__main__":
    build_presentation()
    print(f"Created presentation: {OUTPUT}")
