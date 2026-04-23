from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

from generate_final_presentation import (
    ASSET_DIR,
    DEEP,
    GOLD,
    INK,
    LINE,
    MUTED,
    NAVY,
    OUTPUT,
    PANEL,
    Presentation,
    RGBColor,
    SKY,
    TEAL,
    WHITE,
    add_bg,
    add_bullets,
    add_footer,
    add_textbox,
    add_title,
    build_presentation,
)


ROOT = Path(__file__).resolve().parents[1]
UPDATED_OUTPUT = ROOT / "Trend-Fx_Final_Presentation_Updated.pptx"


def ensure_base_assets():
    required = [
        ASSET_DIR / "hero_forex.png",
        ASSET_DIR / "architecture_diagram.png",
        ASSET_DIR / "erd_diagram.png",
        ASSET_DIR / "frontend_structure.png",
        ASSET_DIR / "project_board.png",
    ]
    if not all(path.exists() for path in required):
        build_presentation()


def add_panel(slide, x, y, w, h, fill=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LINE
    return shape


def build_updated_presentation():
    ensure_base_assets()

    hero = ASSET_DIR / "hero_forex.png"
    architecture = ASSET_DIR / "architecture_diagram.png"
    frontend = ASSET_DIR / "frontend_structure.png"
    board = ASSET_DIR / "project_board.png"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    slide.shapes.add_picture(str(hero), Inches(7.05), Inches(0.45), width=Inches(5.7))
    add_textbox(slide, Inches(0.8), Inches(0.9), Inches(5.5), Inches(0.7), "Trend-Fx", size=30, color=WHITE, bold=True)
    add_textbox(slide, Inches(0.8), Inches(1.55), Inches(5.8), Inches(0.8), "Final Project Presentation", size=24, color=RGBColor(170, 198, 246), bold=True)
    add_textbox(slide, Inches(0.8), Inches(2.35), Inches(5.8), Inches(1.2), "PSD final review deck covering problem, implementation, sprint journey, and future scope.", size=18, color=RGBColor(223, 232, 247))
    add_textbox(slide, Inches(0.8), Inches(5.9), Inches(5.0), Inches(0.35), "Team: Vignesh + Preet", size=17, color=RGBColor(113, 226, 208), bold=True)
    timeline = add_panel(slide, 0.8, 3.65, 5.6, 1.75, fill=DEEP)
    add_textbox(slide, Inches(1.05), Inches(3.88), Inches(5.0), Inches(0.3), "Presentation Timeline (15-18 minutes)", size=18, color=WHITE, bold=True)
    add_bullets(slide, Inches(1.02), Inches(4.22), Inches(5.0), Inches(0.95), [
        "Introduction / Idea",
        "Problem Statement",
        "Solution + Implementation",
        "Sprint Reviews, Challenges, Future Work",
    ], color=RGBColor(232, 239, 251), size=15)
    add_footer(slide, 1)

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_title(slide, "2. Problem Statement", "Understanding forex signals is difficult for beginners because data is fragmented and fast-moving.", dark=False)
    add_panel(slide, 0.75, 1.75, 5.95, 4.95, fill=PANEL)
    add_textbox(slide, Inches(1.0), Inches(2.0), Inches(5.2), Inches(0.4), "Core Challenges", size=20, color=NAVY, bold=True)
    add_bullets(slide, Inches(0.98), Inches(2.4), Inches(5.2), Inches(3.5), [
        "Forex markets move quickly and are difficult to track manually.",
        "Beginners often switch between many tools for price, indicators, and news.",
        "Technical indicators like RSI and MACD are useful, but not easy to interpret without a unified interface.",
        "Many learners need a guided dashboard instead of raw market feeds.",
    ], size=17)
    slide.shapes.add_picture(str(hero), Inches(7.1), Inches(1.95), width=Inches(5.2))
    add_textbox(slide, Inches(7.15), Inches(5.55), Inches(5.0), Inches(0.65), "The project addresses complexity, fragmentation, and learning barriers.", size=17, color=INK, align=PP_ALIGN.CENTER)
    add_footer(slide, 2)

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(246, 248, 252))
    add_title(slide, "3. Idea / Objective", "Build a simplified forex dashboard that combines authentication, market data, and indicators in one product.", dark=False)
    objectives = [
        ("Simplify analysis", "Bring live prices, charts, and indicators into a single interface."),
        ("Support learning", "Help students understand how RSI, MACD, and basic signals relate to price action."),
        ("Provide secure access", "Offer register, login, logout, and password-reset flows for user accounts."),
        ("Show full-stack delivery", "Demonstrate frontend, backend, database, API integration, and team workflow."),
    ]
    positions = [(0.8, 2.0), (6.85, 2.0), (0.8, 4.15), (6.85, 4.15)]
    fills = [RGBColor(236, 244, 255), RGBColor(236, 248, 241), RGBColor(255, 247, 231), RGBColor(246, 241, 255)]
    for (title, desc), (x, y), fill in zip(objectives, positions, fills):
        add_panel(slide, x, y, 5.6, 1.6, fill=fill)
        add_textbox(slide, Inches(x + 0.22), Inches(y + 0.18), Inches(5.0), Inches(0.3), title, size=18, color=NAVY, bold=True)
        add_textbox(slide, Inches(x + 0.22), Inches(y + 0.58), Inches(5.0), Inches(0.7), desc, size=17, color=INK)
    add_footer(slide, 3)

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_title(slide, "4. System Overview", "High-level flow from user interaction to frontend, backend, and external market data APIs.", dark=False)
    slide.shapes.add_picture(str(architecture), Inches(0.7), Inches(1.6), width=Inches(11.9))
    add_footer(slide, 4)

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(246, 248, 252))
    add_title(slide, "5. Technology Stack", "Planned and implemented technologies used across the project.", dark=False)
    stacks = [
        ("Frontend", ["React", "Vite", "React Router DOM", "Recharts"], RGBColor(236, 244, 255)),
        ("Backend", ["FastAPI", "Uvicorn", "SQLAlchemy", "JWT auth"], RGBColor(236, 248, 241)),
        ("Database", ["SQLite test.db", "SQL model design", "Sprint plan referenced MySQL"], RGBColor(255, 247, 231)),
        ("Tools", ["Git", "GitHub", "Branches", "Pull Requests", "PowerShell / VS Code"], RGBColor(246, 241, 255)),
    ]
    for idx, (title, items, fill) in enumerate(stacks):
        x = 0.8 + idx * 3.1
        add_panel(slide, x, 2.0, 2.8, 4.3, fill=fill)
        add_textbox(slide, Inches(x + 0.18), Inches(2.2), Inches(2.3), Inches(0.3), title, size=19, color=NAVY, bold=True)
        add_bullets(slide, Inches(x + 0.15), Inches(2.65), Inches(2.35), Inches(3.1), items, size=15)
    add_footer(slide, 5)

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_title(slide, "6. Implementation", "How the app pulls data, computes indicators, and presents charts to users.", dark=False)
    add_panel(slide, 0.8, 1.85, 4.0, 4.9, fill=PANEL)
    add_textbox(slide, Inches(1.05), Inches(2.1), Inches(3.3), Inches(0.3), "Implementation Highlights", size=20, color=NAVY, bold=True)
    add_bullets(slide, Inches(1.02), Inches(2.5), Inches(3.4), Inches(3.7), [
        "FastAPI backend exposes auth and market endpoints.",
        "Twelve Data API provides live forex values and time-series data.",
        "RSI and MACD are calculated in backend logic before being sent to the dashboard.",
        "Recharts renders line, area, and composed visualizations in the React UI.",
    ], size=17)
    add_panel(slide, 5.15, 1.85, 7.0, 4.9, fill=WHITE)
    add_textbox(slide, Inches(5.45), Inches(2.08), Inches(6.4), Inches(0.3), "Indicator Flow", size=20, color=NAVY, bold=True)
    flow_steps = [
        ("1", "Fetch forex candles", SKY),
        ("2", "Calculate RSI / MACD", TEAL),
        ("3", "Return structured API payload", GOLD),
        ("4", "Render dashboard charts", RGBColor(181, 121, 255)),
    ]
    y = 2.55
    for num, text, color in flow_steps:
        step = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(y), Inches(5.85), Inches(0.72))
        step.fill.solid()
        step.fill.fore_color.rgb = RGBColor(248, 250, 253)
        step.line.color.rgb = LINE
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.7), Inches(y + 0.09), Inches(0.48), Inches(0.48))
        badge.fill.solid()
        badge.fill.fore_color.rgb = color
        badge.line.color.rgb = color
        add_textbox(slide, Inches(5.83), Inches(y + 0.14), Inches(0.2), Inches(0.18), num, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(6.3), Inches(y + 0.15), Inches(4.7), Inches(0.25), text, size=17, color=INK)
        y += 0.9
    add_footer(slide, 6)

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(246, 248, 252))
    add_title(slide, "7. Dashboard Demo", "Live market snapshot, graphs, and quick user-facing analytics.", dark=False)
    slide.shapes.add_picture(str(hero), Inches(0.8), Inches(1.8), width=Inches(7.2))
    add_panel(slide, 8.35, 1.85, 4.2, 4.9, fill=WHITE)
    add_textbox(slide, Inches(8.62), Inches(2.08), Inches(3.6), Inches(0.3), "Demo Walkthrough", size=20, color=NAVY, bold=True)
    add_bullets(slide, Inches(8.58), Inches(2.45), Inches(3.45), Inches(3.8), [
        "User selects a currency pair.",
        "Live price updates appear in the dashboard.",
        "RSI and MACD charts help interpret momentum and trend.",
        "Auth-enabled flow lets users access dashboard after login.",
    ], size=17)
    add_footer(slide, 7)

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_title(slide, "8. Code Explanation", "Recharts powers the dashboard visuals using reusable React components.", dark=False)
    code_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.85), Inches(7.0), Inches(4.9))
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = NAVY
    code_box.line.color.rgb = NAVY
    code_text = """<ResponsiveContainer width="100%" height={280}>
  <LineChart data={chartData}>
    <CartesianGrid strokeDasharray="3 3" />
    <XAxis dataKey="time" />
    <YAxis />
    <Tooltip />
    <Legend />
    <Line dataKey="price" stroke="#67a4ff" />
  </LineChart>
</ResponsiveContainer>"""
    add_textbox(slide, Inches(1.05), Inches(2.15), Inches(6.45), Inches(4.0), code_text, size=16, color=WHITE)
    add_panel(slide, 8.2, 1.85, 4.4, 4.9, fill=PANEL)
    add_textbox(slide, Inches(8.48), Inches(2.08), Inches(3.7), Inches(0.3), "Explanation", size=20, color=NAVY, bold=True)
    add_bullets(slide, Inches(8.42), Inches(2.45), Inches(3.55), Inches(3.8), [
        "ResponsiveContainer makes the chart resize with the layout.",
        "LineChart uses API response data directly from React state.",
        "Tooltip and Legend improve readability during demo and analysis.",
        "The same pattern is reused for RSI and MACD visual sections.",
    ], size=17)
    add_footer(slide, 8)

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(246, 248, 252))
    add_title(slide, "9. Sprint Planning", "Three sprint phases guided the team from setup to analytics and final integration.", dark=False)
    sprint_data = [
        ("Sprint 1", "Project setup, frontend base, layout, routing", RGBColor(236, 244, 255)),
        ("Sprint 2", "Backend APIs, authentication, live data integration", RGBColor(236, 248, 241)),
        ("Sprint 3", "Indicators, dashboard polish, bug fixes, presentation prep", RGBColor(255, 247, 231)),
    ]
    x = 0.9
    for idx, (title, text, fill) in enumerate(sprint_data):
        add_panel(slide, x, 2.35, 3.75, 2.8, fill=fill)
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.2), Inches(2.18), Inches(0.62), Inches(0.62))
        badge.fill.solid()
        badge.fill.fore_color.rgb = [SKY, TEAL, GOLD][idx]
        badge.line.color.rgb = [SKY, TEAL, GOLD][idx]
        add_textbox(slide, Inches(x + 0.33), Inches(2.33), Inches(0.18), Inches(0.18), str(idx + 1), size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(x + 0.22), Inches(2.7), Inches(3.1), Inches(0.3), title, size=19, color=NAVY, bold=True)
        add_textbox(slide, Inches(x + 0.22), Inches(3.2), Inches(3.1), Inches(1.0), text, size=17, color=INK)
        x += 4.1
    add_panel(slide, 1.1, 5.55, 11.1, 0.78, fill=WHITE)
    add_textbox(slide, Inches(1.35), Inches(5.8), Inches(10.6), Inches(0.2), "Planning helped divide ownership clearly between frontend/dashboard work and backend/API/auth work.", size=17, color=INK, align=PP_ALIGN.CENTER)
    add_footer(slide, 9)

    # Slide 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_title(slide, "10. Sprint Review", "Completed work, merged features, and remaining scope after final integration.", dark=False)
    slide.shapes.add_picture(str(board), Inches(0.75), Inches(1.75), width=Inches(7.0))
    add_panel(slide, 8.0, 1.8, 4.55, 4.95, fill=PANEL)
    add_textbox(slide, Inches(8.28), Inches(2.05), Inches(3.8), Inches(0.3), "Status Summary", size=20, color=NAVY, bold=True)
    add_bullets(slide, Inches(8.24), Inches(2.45), Inches(3.65), Inches(3.85), [
        "Completed: auth flows, live market integration, charts, dashboard pages, navigation",
        "Validated: feature branches merged through PRs and review",
        "Pending / future: production email credentials, deeper prediction logic, additional indicators",
    ], size=17)
    add_footer(slide, 10)

    # Slide 11
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(246, 248, 252))
    add_title(slide, "11. Challenges", "Technical and delivery hurdles encountered during the project.", dark=False)
    challenges = [
        "API connectivity and live-data consistency",
        "Indicator calculations and dashboard synchronization",
        "Authentication bugs and route handling",
        "Time management across frontend and backend streams",
    ]
    y = 2.0
    for item in challenges:
        add_panel(slide, 1.0, y, 11.2, 0.88, fill=WHITE)
        add_textbox(slide, Inches(1.28), Inches(y + 0.22), Inches(10.3), Inches(0.22), item, size=19, color=INK)
        y += 1.05
    add_footer(slide, 11)

    # Slide 12
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_title(slide, "12. Solutions", "How the team resolved blockers and stabilized the final build.", dark=False)
    solutions = [
        ("Backend fixes", "Refined FastAPI routing, improved auth flow, and cleaned API integration issues."),
        ("Chart optimization", "Used structured payloads and Recharts components to simplify rendering logic."),
        ("Iteration via branches", "Separated work into feature branches and merged after validation."),
        ("Focused sprint reviews", "Used sprint checkpoints to identify bugs early and keep the scope realistic."),
    ]
    positions = [(0.85, 2.0), (6.7, 2.0), (0.85, 4.25), (6.7, 4.25)]
    fills = [RGBColor(236, 248, 241), RGBColor(236, 244, 255), RGBColor(246, 241, 255), RGBColor(255, 247, 231)]
    for (title, desc), (x, y), fill in zip(solutions, positions, fills):
        add_panel(slide, x, y, 5.6, 1.7, fill=fill)
        add_textbox(slide, Inches(x + 0.22), Inches(y + 0.18), Inches(4.9), Inches(0.28), title, size=18, color=NAVY, bold=True)
        add_textbox(slide, Inches(x + 0.22), Inches(y + 0.58), Inches(4.9), Inches(0.78), desc, size=17, color=INK)
    add_footer(slide, 12)

    # Slide 13
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, RGBColor(246, 248, 252))
    add_title(slide, "13. Future Work", "Planned enhancements after the current academic version.", dark=False)
    add_bullets(slide, Inches(0.9), Inches(2.0), Inches(5.8), Inches(3.9), [
        "Add more technical indicators and stronger signal explanations",
        "Improve prediction logic and confidence scoring",
        "Introduce richer news sentiment and backtesting support",
        "Move toward production-ready email delivery and deployment",
    ], size=19)
    add_panel(slide, 7.25, 2.0, 5.0, 3.8, fill=WHITE)
    add_textbox(slide, Inches(7.55), Inches(2.35), Inches(4.3), Inches(0.35), "Long-Term Vision", size=21, color=NAVY, bold=True)
    add_textbox(slide, Inches(7.55), Inches(2.9), Inches(4.2), Inches(1.8), "Evolve Trend-Fx from a learning-focused dashboard into a stronger decision-support platform with deeper analytics and more robust forecasting.", size=19, color=INK)
    add_footer(slide, 13)

    # Slide 14
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, WHITE)
    add_title(slide, "14. Conclusion", "Trend-Fx demonstrates a complete full-stack solution to simplify forex learning and analysis.", dark=False)
    add_panel(slide, 0.9, 2.0, 11.5, 3.8, fill=PANEL)
    add_bullets(slide, Inches(1.2), Inches(2.45), Inches(10.8), Inches(2.7), [
        "We identified a real usability problem in forex learning and built a focused dashboard solution.",
        "The project integrates React UI, FastAPI services, authentication, database persistence, and live market data.",
        "Sprint-based teamwork and repository workflow helped deliver the final result with clear ownership.",
        "The platform now provides a strong academic foundation for future prediction improvements.",
    ], size=19)
    add_footer(slide, 14)

    # Slide 15
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, NAVY)
    add_textbox(slide, Inches(1.1), Inches(2.0), Inches(11.0), Inches(0.8), "Thank You", size=30, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.1), Inches(2.9), Inches(11.0), Inches(0.55), "Questions?", size=24, color=RGBColor(112, 226, 208), bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(2.0), Inches(4.1), Inches(9.2), Inches(0.7), "Trend-Fx Final Presentation", size=18, color=RGBColor(214, 226, 246), align=PP_ALIGN.CENTER)
    add_footer(slide, 15)

    prs.save(UPDATED_OUTPUT)
    print(f"Created updated presentation: {UPDATED_OUTPUT}")


if __name__ == "__main__":
    build_updated_presentation()
